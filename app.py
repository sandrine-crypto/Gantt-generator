"""
GANTT CHART GENERATOR - Application Streamlit
Interface web pour generer des diagrammes de Gantt a partir de fichiers Excel.
"""

import io
from dataclasses import dataclass, field
from datetime import datetime
from typing import TYPE_CHECKING, List

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM
from PIL import Image

if TYPE_CHECKING:
    from pandas import DataFrame, Timestamp


# =============================================================================
# CONFIGURATION
# =============================================================================

@dataclass
class GanttConfig:
    """Configuration centralisee pour les Gantt charts."""
    BAR_HEIGHT: int = 35
    LEFT_MARGIN: int = 260
    TOP_MARGIN: int = 100
    TOTAL_WIDTH: int = 1300
    RIGHT_MARGIN: int = 50
    BOTTOM_MARGIN: int = 100
    MIN_BAR_WIDTH: int = 3
    MIN_BAR_WIDTH_FOR_TEXT: int = 40
    GRID_STEP_DAYS: int = 90
    MAX_CODE_LENGTH: int = 25
    MAX_TARGET_LENGTH: int = 25


# Couleurs par defaut
DEFAULT_COLORS = [
    '#3498db',  # Bleu
    '#e74c3c',  # Rouge
    '#2ecc71',  # Vert
    '#f39c12',  # Orange
    '#9b59b6',  # Violet
    '#1abc9c',  # Turquoise
    '#e67e22',  # Orange fonce
    '#34495e',  # Gris fonce
    '#16a085',  # Vert fonce
    '#c0392b'   # Rouge fonce
]

CONFIG = GanttConfig()

SVG_STYLES = """
.gantt-label { font-size: 12px; font-family: Arial; }
.gantt-label-code { font-weight: bold; font-size: 13px; fill: #2c3e50; }
.gantt-label-target { font-size: 10px; fill: #7f8c8d; }
.gantt-title { font-size: 18px; font-weight: bold; fill: #2c3e50; }
.gantt-bar { opacity: 0.85; stroke: white; stroke-width: 1; }
.gantt-bar-text { font-size: 11px; fill: white; font-weight: bold; text-anchor: middle; }
.gantt-grid-line { stroke: #ecf0f1; stroke-width: 1; }
.gantt-date { font-size: 10px; fill: #7f8c8d; text-anchor: middle; }
.gantt-legend { font-size: 9px; fill: #7f8c8d; }
""".strip()

HTML_STYLES = """
* { margin: 0; padding: 0; box-sizing: border-box; }
body {
    font-family: "Segoe UI", Tahoma, Geneva, Verdana, sans-serif;
    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    min-height: 100vh;
    padding: 20px;
}
.slide-container {
    background: white;
    border-radius: 10px;
    box-shadow: 0 10px 40px rgba(0,0,0,0.2);
    margin: 20px auto;
    padding: 30px;
    max-width: 1400px;
    page-break-after: always;
}
.slide-header {
    border-bottom: 3px solid #667eea;
    padding-bottom: 15px;
    margin-bottom: 20px;
}
.slide-title { font-size: 28px; color: #2c3e50; font-weight: bold; margin-bottom: 5px; }
.slide-subtitle { font-size: 14px; color: #7f8c8d; }
.gantt-wrapper { overflow-x: auto; margin-top: 20px; }
.gantt-wrapper svg { display: block; margin: 0 auto; }
.summary-table { width: 100%; border-collapse: collapse; font-size: 14px; }
.summary-table th, .summary-table td { border: 1px solid #bdc3c7; padding: 12px; }
.summary-table th { background: #ecf0f1; text-align: left; }
.summary-table td:last-child { text-align: center; }
@media print { .slide-container { page-break-after: always; margin: 0; } }
""".strip()


# =============================================================================
# CLASSE GENERATEUR
# =============================================================================

class GanttGenerator:
    """Generateur de diagrammes de Gantt pour modeles murins."""

    def __init__(self, df: "DataFrame", colors: List[str] = None) -> None:
        self.df = df
        self.df_clean: "DataFrame" = None
        self.gantt_svgs: dict = {}
        self.min_date: "Timestamp" = None
        self.max_date: "Timestamp" = None
        self._date_range_days: int = 0
        self.colors = colors if colors else DEFAULT_COLORS.copy()

    def load_data(self) -> "GanttGenerator":
        """Prepare les donnees."""
        col_start = 'date disponibilité HO'
        col_end = 'data de fin de validation (fin du dernier MI critique taggé validation)'

        df = self.df.copy()
        df['internal code'] = df['internal code'].fillna('N/A')
        df['target'] = df['target'].fillna('N/A')
        df[col_start] = pd.to_datetime(df[col_start], errors='coerce')
        df[col_end] = pd.to_datetime(df[col_end], errors='coerce')

        mask = (
            (df['internal code'] != 'N/A') &
            df[col_start].notna() &
            df[col_end].notna()
        )
        self.df_clean = df.loc[mask].copy()

        self.df_clean['start_date'] = self.df_clean[col_start]
        self.df_clean['end_date'] = self.df_clean[col_end]
        self.df_clean['duration_days'] = (
            self.df_clean['end_date'] - self.df_clean['start_date']
        ).dt.days

        self.min_date = self.df_clean['start_date'].min()
        self.max_date = self.df_clean['end_date'].max()
        self._date_range_days = (self.max_date - self.min_date).days

        return self

    def _compute_x_position(self, date: "Timestamp", chart_width: float) -> float:
        days_from_start = (date - self.min_date).days
        return CONFIG.LEFT_MARGIN + (days_from_start / self._date_range_days) * chart_width

    def _build_svg_header(self, total_width: int, total_height: int,
                          product_line_name: str, model_count: int) -> list:
        return [
            f'<svg width="{total_width}" height="{total_height}" xmlns="http://www.w3.org/2000/svg">',
            f'<defs><style>{SVG_STYLES}</style></defs>',
            f'<rect width="{total_width}" height="{total_height}" fill="white"/>',
            f'<text x="{CONFIG.LEFT_MARGIN}" y="35" class="gantt-title">{product_line_name}</text>',
            f'<text x="{CONFIG.LEFT_MARGIN}" y="55" class="gantt-legend">({model_count} modeles)</text>',
        ]

    def _build_date_grid(self, chart_width: float, total_height: int) -> list:
        parts = []
        date_range = pd.date_range(
            start=self.min_date,
            end=self.max_date,
            freq=f'{CONFIG.GRID_STEP_DAYS}D'
        )

        for current in date_range:
            x_pos = self._compute_x_position(current, chart_width)
            date_str = current.strftime('%b %Y')
            parts.append(
                f'<line x1="{x_pos}" y1="{CONFIG.TOP_MARGIN}" '
                f'x2="{x_pos}" y2="{total_height - 50}" class="gantt-grid-line"/>'
            )
            parts.append(
                f'<text x="{x_pos}" y="{CONFIG.TOP_MARGIN - 15}" class="gantt-date">{date_str}</text>'
            )
        return parts

    def _build_gantt_bar(self, row: pd.Series, idx: int, chart_width: float) -> list:
        parts = []
        y_pos = CONFIG.TOP_MARGIN + idx * CONFIG.BAR_HEIGHT + 15

        code = str(row['internal code'])[:CONFIG.MAX_CODE_LENGTH]
        target = str(row['target'])[:CONFIG.MAX_TARGET_LENGTH]

        parts.append(f'<text x="10" y="{y_pos + 5}" class="gantt-label gantt-label-code">{code}</text>')
        parts.append(f'<text x="10" y="{y_pos + 18}" class="gantt-label gantt-label-target">{target}</text>')

        start_x = self._compute_x_position(row['start_date'], chart_width)
        end_x = self._compute_x_position(row['end_date'], chart_width)
        bar_width = max(end_x - start_x, CONFIG.MIN_BAR_WIDTH)
        color = self.colors[idx % len(self.colors)]

        ho_date = row['start_date'].strftime('%d/%m/%Y')
        val_date = row['end_date'].strftime('%d/%m/%Y')
        duration = int(row['duration_days'])

        tooltip = f"{code} | {target} | HO: {ho_date} -> Val: {val_date} | {duration}j"
        parts.append(
            f'<rect x="{start_x}" y="{y_pos - 8}" width="{bar_width}" height="25" '
            f'fill="{color}" class="gantt-bar"><title>{tooltip}</title></rect>'
        )

        if bar_width > CONFIG.MIN_BAR_WIDTH_FOR_TEXT:
            text_x = start_x + bar_width / 2
            parts.append(f'<text x="{text_x}" y="{y_pos + 2}" class="gantt-bar-text">{duration}j</text>')

        return parts

    def create_gantt_svg(self, product_line_data: "DataFrame", product_line_name: str) -> str:
        data = product_line_data.sort_values('start_date').reset_index(drop=True)
        model_count = len(data)

        chart_width = CONFIG.TOTAL_WIDTH - CONFIG.LEFT_MARGIN - CONFIG.RIGHT_MARGIN
        total_height = model_count * CONFIG.BAR_HEIGHT + CONFIG.TOP_MARGIN + CONFIG.BOTTOM_MARGIN

        svg_parts = []
        svg_parts.extend(self._build_svg_header(CONFIG.TOTAL_WIDTH, total_height, product_line_name, model_count))
        svg_parts.extend(self._build_date_grid(chart_width, total_height))

        for idx, row in data.iterrows():
            svg_parts.extend(self._build_gantt_bar(row, idx, chart_width))

        axis_y = CONFIG.TOP_MARGIN + model_count * CONFIG.BAR_HEIGHT
        svg_parts.append(
            f'<line x1="{CONFIG.LEFT_MARGIN}" y1="{axis_y}" '
            f'x2="{CONFIG.TOTAL_WIDTH - CONFIG.RIGHT_MARGIN}" y2="{axis_y}" '
            f'stroke="#34495e" stroke-width="2"/>'
        )
        svg_parts.append(
            f'<text x="{CONFIG.LEFT_MARGIN}" y="{total_height - 30}" class="gantt-legend">'
            f'Disponibilite HO -> Fin Validation</text>'
        )
        svg_parts.append('</svg>')

        return ''.join(svg_parts)

    def generate_gantt_charts(self) -> "GanttGenerator":
        grouped = self.df_clean.groupby('product line', sort=True)
        for pl_name, pl_data in grouped:
            self.gantt_svgs[pl_name] = self.create_gantt_svg(pl_data, pl_name)
        return self

    def generate_html_slides(self) -> str:
        html_parts = [
            '<!DOCTYPE html>',
            '<html lang="fr">',
            '<head>',
            '<meta charset="UTF-8">',
            '<meta name="viewport" content="width=device-width, initial-scale=1.0">',
            '<title>Gantt Charts - Modeles Murins CRUPPE</title>',
            f'<style>{HTML_STYLES}</style>',
            '</head>',
            '<body>',
        ]

        # Cover slide
        html_parts.append(f"""
<div class="slide-container">
    <div class="slide-header">
        <div class="slide-title">Diagrammes de Gantt</div>
        <div class="slide-subtitle">Modeles Murins Genetiquement Modifies - CRUPPE</div>
    </div>
    <div style="text-align: center; margin-top: 60px;">
        <p style="font-size: 18px; color: #7f8c8d; margin: 20px 0;">
            <strong>{len(self.df_clean)} modeles</strong> repartis dans
            <strong>{len(self.gantt_svgs)} product lines</strong>
        </p>
        <p style="font-size: 14px; color: #95a5a6; margin: 20px 0;">
            Periode: {self.min_date.strftime('%d/%m/%Y')} -> {self.max_date.strftime('%d/%m/%Y')}
        </p>
        <p style="font-size: 12px; color: #bdc3c7; margin: 40px 0;">
            Genere le {datetime.now().strftime('%d/%m/%Y a %H:%M')}
        </p>
    </div>
</div>""")

        # Gantt slides
        for idx, (product_line, svg) in enumerate(self.gantt_svgs.items(), 1):
            html_parts.append(f"""
<div class="slide-container">
    <div class="slide-header">
        <div class="slide-title">Produit {idx}/{len(self.gantt_svgs)}</div>
        <div class="slide-subtitle">{product_line}</div>
    </div>
    <div class="gantt-wrapper">{svg}</div>
</div>""")

        # Summary slide
        counts = self.df_clean['product line'].value_counts().sort_values(ascending=False)
        rows = ''.join(f'<tr><td>{pl}</td><td>{count}</td></tr>' for pl, count in counts.items())
        html_parts.append(f"""
<div class="slide-container">
    <div class="slide-header">
        <div class="slide-title">Resume</div>
        <div class="slide-subtitle">Vue d'ensemble des donnees</div>
    </div>
    <div style="margin-top: 40px;">
        <table class="summary-table">
            <tr><th>Product Line</th><th>Modeles</th></tr>
            {rows}
        </table>
    </div>
</div>""")

        html_parts.extend(['</body>', '</html>'])
        return ''.join(html_parts)

    def export_csv(self) -> str:
        column_mapping = {
            'product line': 'Product Line',
            'internal code': 'Internal Code',
            'target': 'Target',
            'status': 'Status',
            'start_date': 'Date Disponibilite HO',
            'end_date': 'Date Fin Validation',
            'duration_days': 'Duree (jours)',
        }

        available_cols = [c for c in column_mapping.keys() if c in self.df_clean.columns]
        export_df = (
            self.df_clean[available_cols]
            .rename(columns=column_mapping)
            .sort_values('Product Line')
        )

        return export_df.to_csv(index=False)

    def _svg_to_png_bytes(self, svg_content: str) -> bytes:
        """Convertit un SVG en bytes PNG."""
        svg_io = io.BytesIO(svg_content.encode('utf-8'))
        drawing = svg2rlg(svg_io)
        png_io = io.BytesIO()
        renderPM.drawToFile(drawing, png_io, fmt='PNG', dpi=150)
        png_io.seek(0)
        return png_io.getvalue()

    def export_pptx(self) -> bytes:
        """Genere une presentation PowerPoint avec les Gantt charts."""
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide de couverture
        blank_layout = prs.slide_layouts[6]  # Layout vide
        slide = prs.slides.add_slide(blank_layout)

        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Diagrammes de Gantt"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Sous-titre
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.33), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = "Modeles Murins Genetiquement Modifies - CRUPPE"
        subtitle_para.font.size = Pt(24)
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Statistiques
        stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.33), Inches(1))
        stats_frame = stats_box.text_frame
        stats_para = stats_frame.paragraphs[0]
        stats_para.text = f"{len(self.df_clean)} modeles | {len(self.gantt_svgs)} product lines"
        stats_para.font.size = Pt(18)
        stats_para.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.33), Inches(0.5))
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = f"Periode: {self.min_date.strftime('%d/%m/%Y')} - {self.max_date.strftime('%d/%m/%Y')}"
        date_para.font.size = Pt(14)
        date_para.alignment = PP_ALIGN.CENTER

        # Slides pour chaque product line
        for idx, (product_line, svg) in enumerate(self.gantt_svgs.items(), 1):
            slide = prs.slides.add_slide(blank_layout)

            # Titre de la slide
            title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.6))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"{product_line} ({idx}/{len(self.gantt_svgs)})"
            title_para.font.size = Pt(24)
            title_para.font.bold = True

            # Convertir SVG en PNG et l'ajouter
            try:
                png_bytes = self._svg_to_png_bytes(svg)
                png_io = io.BytesIO(png_bytes)

                # Calculer les dimensions pour adapter l'image
                with Image.open(io.BytesIO(png_bytes)) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height

                    max_width = 12.5
                    max_height = 6.5
                    if aspect_ratio > (max_width / max_height):
                        width = Inches(max_width)
                        height = Inches(max_width / aspect_ratio)
                    else:
                        height = Inches(max_height)
                        width = Inches(max_height * aspect_ratio)

                left = Inches((13.33 - width.inches) / 2)
                top = Inches(0.9)
                slide.shapes.add_picture(png_io, left, top, width=width, height=height)
            except Exception:
                # En cas d'erreur, ajouter un texte
                error_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
                error_frame = error_box.text_frame
                error_para = error_frame.paragraphs[0]
                error_para.text = f"Erreur lors de la generation du graphique pour {product_line}"
                error_para.alignment = PP_ALIGN.CENTER

        # Slide de resume
        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Resume"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Tableau de resume
        counts = self.df_clean['product line'].value_counts().sort_values(ascending=False)
        rows = len(counts) + 1
        cols = 2

        table_width = Inches(8)
        table_height = Inches(min(rows * 0.4, 5.5))
        left = Inches((13.33 - 8) / 2)
        top = Inches(1.2)

        table = slide.shapes.add_table(rows, cols, left, top, table_width, table_height).table

        # En-tetes
        table.cell(0, 0).text = "Product Line"
        table.cell(0, 1).text = "Modeles"

        for row_idx, (pl, count) in enumerate(counts.items(), 1):
            table.cell(row_idx, 0).text = str(pl)
            table.cell(row_idx, 1).text = str(count)

        # Sauvegarder
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io.getvalue()

    def export_docx(self) -> bytes:
        """Genere un document Word avec les Gantt charts."""
        doc = Document()

        # Titre
        title = doc.add_heading('Diagrammes de Gantt', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Sous-titre
        subtitle = doc.add_paragraph('Modeles Murins Genetiquement Modifies - CRUPPE')
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Statistiques
        stats = doc.add_paragraph()
        stats.alignment = WD_ALIGN_PARAGRAPH.CENTER
        stats.add_run(f"{len(self.df_clean)} modeles").bold = True
        stats.add_run(f" repartis dans ")
        stats.add_run(f"{len(self.gantt_svgs)} product lines").bold = True

        doc.add_paragraph(
            f"Periode: {self.min_date.strftime('%d/%m/%Y')} - {self.max_date.strftime('%d/%m/%Y')}"
        ).alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph(
            f"Genere le {datetime.now().strftime('%d/%m/%Y a %H:%M')}"
        ).alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_page_break()

        # Pages pour chaque product line
        for idx, (product_line, svg) in enumerate(self.gantt_svgs.items(), 1):
            doc.add_heading(f'{product_line} ({idx}/{len(self.gantt_svgs)})', level=1)

            # Convertir SVG en PNG et l'ajouter
            try:
                png_bytes = self._svg_to_png_bytes(svg)
                png_io = io.BytesIO(png_bytes)

                # Calculer la largeur optimale
                with Image.open(io.BytesIO(png_bytes)) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    doc_width = 6.5  # pouces
                    doc_height = doc_width / aspect_ratio

                doc.add_picture(png_io, width=DocxInches(doc_width))
            except Exception:
                doc.add_paragraph(f"Erreur lors de la generation du graphique pour {product_line}")

            doc.add_page_break()

        # Page de resume
        doc.add_heading('Resume', level=1)

        # Tableau de resume
        counts = self.df_clean['product line'].value_counts().sort_values(ascending=False)
        table = doc.add_table(rows=len(counts) + 1, cols=2)
        table.style = 'Table Grid'

        # En-tetes
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Product Line'
        hdr_cells[1].text = 'Modeles'
        for cell in hdr_cells:
            cell.paragraphs[0].runs[0].bold = True

        # Donnees
        for row_idx, (pl, count) in enumerate(counts.items(), 1):
            row_cells = table.rows[row_idx].cells
            row_cells[0].text = str(pl)
            row_cells[1].text = str(count)

        # Sauvegarder
        docx_io = io.BytesIO()
        doc.save(docx_io)
        docx_io.seek(0)
        return docx_io.getvalue()


# =============================================================================
# APPLICATION STREAMLIT
# =============================================================================

def main():
    st.set_page_config(
        page_title="Gantt Chart Generator - CRUPPE",
        page_icon="📊",
        layout="wide"
    )

    st.title("📊 Generateur de Gantt Charts")
    st.markdown("**Modeles Murins Genetiquement Modifies - CRUPPE**")
    st.markdown("---")

    # Instructions
    with st.expander("📋 Instructions d'utilisation"):
        st.markdown("""
        ### Comment utiliser cette application:

        1. **Preparez votre fichier Excel** avec une feuille nommee `tableau complet`
        2. **Colonnes requises:**
           - `internal code` - Code unique du modele
           - `target` - Gene/proteine ciblee
           - `product line` - Categorie du modele
           - `date disponibilite HO` - Date de debut
           - `data de fin de validation (fin du dernier MI critique tagge validation)` - Date de fin
        3. **Personnalisez les couleurs** des barres du Gantt (optionnel)
        4. **Uploadez le fichier** via le bouton ci-dessous
        5. **Telechargez les resultats** (PowerPoint, Word ou CSV)
        """)

    # File upload
    st.subheader("1. Charger le fichier Excel")
    uploaded_file = st.file_uploader(
        "Selectionnez votre fichier Excel",
        type=['xlsx', 'xls'],
        help="Le fichier doit contenir une feuille 'tableau complet'"
    )

    # Color customization section
    st.subheader("2. Personnaliser les couleurs")
    with st.expander("Modifier les couleurs des barres", expanded=False):
        st.markdown("Cliquez sur chaque couleur pour la modifier:")
        color_cols = st.columns(5)
        custom_colors = []
        color_names = [
            "Bleu", "Rouge", "Vert", "Orange", "Violet",
            "Turquoise", "Orange fonce", "Gris fonce", "Vert fonce", "Rouge fonce"
        ]
        for i, default_color in enumerate(DEFAULT_COLORS):
            col_idx = i % 5
            with color_cols[col_idx]:
                color = st.color_picker(
                    color_names[i],
                    default_color,
                    key=f"color_{i}"
                )
                custom_colors.append(color)

    if uploaded_file is not None:
        try:
            # Load data
            with st.spinner("Chargement des donnees..."):
                df = pd.read_excel(uploaded_file, sheet_name='tableau complet')

            st.success(f"Fichier charge: {len(df)} lignes trouvees")

            # Generate charts
            st.subheader("3. Generation des Gantt Charts")

            with st.spinner("Generation en cours..."):
                generator = GanttGenerator(df, colors=custom_colors)
                generator.load_data()
                generator.generate_gantt_charts()

            # Statistics
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Modeles valides", len(generator.df_clean))
            with col2:
                st.metric("Product Lines", len(generator.gantt_svgs))
            with col3:
                if generator._date_range_days > 0:
                    st.metric("Periode", f"{generator._date_range_days} jours")

            # Display charts
            st.subheader("4. Visualisation des Gantt Charts")

            selected_pl = st.selectbox(
                "Selectionnez une Product Line:",
                options=list(generator.gantt_svgs.keys())
            )

            if selected_pl:
                st.markdown(f"**{selected_pl}**")
                # Display SVG
                st.components.v1.html(
                    generator.gantt_svgs[selected_pl],
                    height=len(generator.df_clean[generator.df_clean['product line'] == selected_pl]) * 40 + 200,
                    scrolling=True
                )

            # Downloads
            st.subheader("5. Telecharger les resultats")

            col1, col2, col3 = st.columns(3)

            with col1:
                with st.spinner("Preparation PowerPoint..."):
                    pptx_content = generator.export_pptx()
                st.download_button(
                    label="📽️ Telecharger PowerPoint (.pptx)",
                    data=pptx_content,
                    file_name="gantt_slides.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            with col2:
                with st.spinner("Preparation Word..."):
                    docx_content = generator.export_docx()
                st.download_button(
                    label="📄 Telecharger Word (.docx)",
                    data=docx_content,
                    file_name="gantt_document.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )

            with col3:
                csv_content = generator.export_csv()
                st.download_button(
                    label="📊 Telecharger CSV (donnees)",
                    data=csv_content,
                    file_name="gantt_models_export.csv",
                    mime="text/csv"
                )

            # Summary table
            st.subheader("6. Resume par Product Line")
            summary = generator.df_clean['product line'].value_counts().reset_index()
            summary.columns = ['Product Line', 'Nombre de modeles']
            st.dataframe(summary, use_container_width=True)

        except ValueError as e:
            st.error(f"Erreur: La feuille 'tableau complet' n'a pas ete trouvee dans le fichier Excel.")
            st.info("Verifiez que votre fichier contient bien une feuille nommee exactement 'tableau complet'")
        except Exception as e:
            st.error(f"Erreur lors du traitement: {str(e)}")

    else:
        st.info("👆 Veuillez uploader un fichier Excel pour commencer")

    # Footer
    st.markdown("---")
    st.markdown(
        "<div style='text-align: center; color: gray;'>"
        "CRUPPE - Biologie Moleculaire | Lyon, Rhone-Alpes"
        "</div>",
        unsafe_allow_html=True
    )


if __name__ == "__main__":
    main()
