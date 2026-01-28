#!/usr/bin/env python3
"""
GANTT CHART GENERATOR - Application Streamlit
Génère des diagrammes de Gantt interactifs avec export HTML, CSV, PPTX et DOCX
Couleurs personnalisables pour les diagrammes

USAGE:
    streamlit run gantt_generator_V2.py

REQUIREMENTS:
    streamlit, pandas, openpyxl, python-pptx, python-docx, matplotlib, Pillow
"""

from __future__ import annotations

import io
import base64
from dataclasses import dataclass, field
from datetime import datetime
from typing import TYPE_CHECKING

import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.figure import Figure

# Pour export PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN

# Pour export DOCX
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

if TYPE_CHECKING:
    from pandas import DataFrame, Timestamp


# =============================================================================
# CONSTANTES DE CONFIGURATION
# =============================================================================

@dataclass
class GanttConfig:
    """Configuration centralisée pour les Gantt charts."""

    # Dimensions SVG
    BAR_HEIGHT: int = 35
    LEFT_MARGIN: int = 260
    TOP_MARGIN: int = 100
    TOTAL_WIDTH: int = 1300
    RIGHT_MARGIN: int = 50
    BOTTOM_MARGIN: int = 100
    MIN_BAR_WIDTH: int = 3
    MIN_BAR_WIDTH_FOR_TEXT: int = 40

    # Grille temporelle
    GRID_STEP_DAYS: int = 90

    # Palette de couleurs par défaut
    DEFAULT_COLORS: tuple[str, ...] = (
        '#3498db', '#e74c3c', '#2ecc71', '#f39c12', '#9b59b6',
        '#1abc9c', '#e67e22', '#34495e', '#16a085', '#c0392b'
    )

    # Limites de texte
    MAX_CODE_LENGTH: int = 25
    MAX_TARGET_LENGTH: int = 25


CONFIG = GanttConfig()

# Styles CSS réutilisables
SVG_STYLES = """
.gantt-label { font-size: 12px; font-family: Arial; }
.gantt-label-code { font-weight: bold; font-size: 13px; fill: #2c3e50; }
.gantt-label-target { font-size: 10px; fill: #7f8c8d; }
.gantt-title { font-size: 18px; font-weight: bold; fill: #2c3e50; }
.gantt-bar { opacity: 0.85; stroke: white; stroke-width: 1; }
.gantt-bar-text { font-size: 11px; fill: white; font-weight: bold; text-anchor: middle; }
.gantt-grid-line { stroke: #ecf0f1; stroke-width: 1; }
.gantt-date { font-size: 10px; fill: #7f8c8d; text-anchor: middle; }
.gantt-legend { font-size: 9px; fill: #7f8c8d; }
""".strip()

HTML_STYLES = """
* { margin: 0; padding: 0; box-sizing: border-box; }
body {
    font-family: "Segoe UI", Tahoma, Geneva, Verdana, sans-serif;
    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    min-height: 100vh;
    padding: 20px;
}
.slide-container {
    background: white;
    border-radius: 10px;
    box-shadow: 0 10px 40px rgba(0,0,0,0.2);
    margin: 20px auto;
    padding: 30px;
    max-width: 1400px;
    page-break-after: always;
}
.slide-header {
    border-bottom: 3px solid #667eea;
    padding-bottom: 15px;
    margin-bottom: 20px;
}
.slide-title { font-size: 28px; color: #2c3e50; font-weight: bold; margin-bottom: 5px; }
.slide-subtitle { font-size: 14px; color: #7f8c8d; }
.gantt-wrapper { overflow-x: auto; margin-top: 20px; }
.gantt-wrapper svg { display: block; margin: 0 auto; }
.summary-table { width: 100%; border-collapse: collapse; font-size: 14px; }
.summary-table th, .summary-table td { border: 1px solid #bdc3c7; padding: 12px; }
.summary-table th { background: #ecf0f1; text-align: left; }
.summary-table td:last-child { text-align: center; }
@media print { .slide-container { page-break-after: always; margin: 0; } }
""".strip()


# =============================================================================
# FONCTIONS UTILITAIRES
# =============================================================================

def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convertit une couleur hexadécimale en RGB."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_matplotlib_gantt(data: DataFrame, product_line: str, colors: list[str]) -> Figure:
    """Crée un diagramme de Gantt avec matplotlib pour l'export."""
    fig, ax = plt.subplots(figsize=(14, max(4, len(data) * 0.5)))

    data_sorted = data.sort_values('start_date').reset_index(drop=True)

    for idx, (_, row) in enumerate(data_sorted.iterrows()):
        color = colors[idx % len(colors)]
        start = row['start_date']
        duration = row['duration_days']

        ax.barh(idx, duration, left=start, height=0.6, color=color, alpha=0.85, edgecolor='white')

        # Ajouter le texte du code
        code = str(row['internal code'])[:CONFIG.MAX_CODE_LENGTH]
        ax.text(start, idx, f"  {code}", va='center', ha='left', fontsize=8, fontweight='bold')

    # Configuration des axes
    ax.set_yticks(range(len(data_sorted)))
    ax.set_yticklabels([str(row['target'])[:CONFIG.MAX_TARGET_LENGTH] for _, row in data_sorted.iterrows()], fontsize=8)
    ax.set_xlabel('Date')
    ax.set_title(f'{product_line} ({len(data)} modeles)', fontweight='bold', fontsize=12)
    ax.invert_yaxis()
    ax.grid(axis='x', alpha=0.3)

    plt.tight_layout()
    return fig


def fig_to_image_bytes(fig: Figure, format: str = 'png', dpi: int = 150) -> bytes:
    """Convertit une figure matplotlib en bytes."""
    buf = io.BytesIO()
    fig.savefig(buf, format=format, dpi=dpi, bbox_inches='tight')
    buf.seek(0)
    plt.close(fig)
    return buf.getvalue()


# =============================================================================
# CLASSE PRINCIPALE
# =============================================================================

class GanttGenerator:
    """Générateur de diagrammes de Gantt pour modèles murins."""

    def __init__(self, df: DataFrame, colors: list[str] = None) -> None:
        """Initialise avec un DataFrame et des couleurs optionnelles."""
        self.df = df
        self.df_clean: DataFrame | None = None
        self.gantt_svgs: dict[str, str] = {}
        self.min_date: Timestamp | None = None
        self.max_date: Timestamp | None = None
        self._date_range_days: int = 0
        self.colors = colors or list(CONFIG.DEFAULT_COLORS)

    def load_data(self) -> 'GanttGenerator':
        """Charge et prépare les données."""
        # Colonnes source
        col_start = 'date disponibilité HO'
        col_end = 'data de fin de validation (fin du dernier MI critique taggé validation)'

        # Nettoyage vectorisé des données
        df = self.df.copy()
        df['internal code'] = df['internal code'].fillna('N/A')
        df['target'] = df['target'].fillna('N/A')
        df[col_start] = pd.to_datetime(df[col_start], errors='coerce')
        df[col_end] = pd.to_datetime(df[col_end], errors='coerce')

        # Filtrage avec masque booléen
        mask = (
            (df['internal code'] != 'N/A') &
            df[col_start].notna() &
            df[col_end].notna()
        )
        self.df_clean = df.loc[mask].copy()

        # Colonnes de travail
        self.df_clean['start_date'] = self.df_clean[col_start]
        self.df_clean['end_date'] = self.df_clean[col_end]
        self.df_clean['duration_days'] = (
            self.df_clean['end_date'] - self.df_clean['start_date']
        ).dt.days

        # Calcul unique des bornes temporelles
        self.min_date = self.df_clean['start_date'].min()
        self.max_date = self.df_clean['end_date'].max()
        self._date_range_days = (self.max_date - self.min_date).days

        return self

    def _compute_x_position(self, date, chart_width: float) -> float:
        """Calcule la position X pour une date donnée."""
        days_from_start = (date - self.min_date).days
        return CONFIG.LEFT_MARGIN + (days_from_start / self._date_range_days) * chart_width

    def _build_svg_header(self, total_width: int, total_height: int,
                          product_line_name: str, model_count: int) -> list[str]:
        """Construit l'en-tête SVG avec styles et titre."""
        return [
            f'<svg width="{total_width}" height="{total_height}" xmlns="http://www.w3.org/2000/svg">',
            f'<defs><style>{SVG_STYLES}</style></defs>',
            f'<rect width="{total_width}" height="{total_height}" fill="white"/>',
            f'<text x="{CONFIG.LEFT_MARGIN}" y="35" class="gantt-title">{product_line_name}</text>',
            f'<text x="{CONFIG.LEFT_MARGIN}" y="55" class="gantt-legend">({model_count} modeles)</text>',
        ]

    def _build_date_grid(self, chart_width: float, total_height: int) -> list[str]:
        """Construit la grille de dates."""
        parts: list[str] = []

        date_range = pd.date_range(
            start=self.min_date,
            end=self.max_date,
            freq=f'{CONFIG.GRID_STEP_DAYS}D'
        )

        for current in date_range:
            x_pos = self._compute_x_position(current, chart_width)
            date_str = current.strftime('%b %Y')
            parts.append(
                f'<line x1="{x_pos}" y1="{CONFIG.TOP_MARGIN}" '
                f'x2="{x_pos}" y2="{total_height - 50}" class="gantt-grid-line"/>'
            )
            parts.append(
                f'<text x="{x_pos}" y="{CONFIG.TOP_MARGIN - 15}" class="gantt-date">{date_str}</text>'
            )

        return parts

    def _build_gantt_bar(self, row: pd.Series, idx: int, chart_width: float) -> list[str]:
        """Construit une barre Gantt avec son label."""
        parts: list[str] = []
        y_pos = CONFIG.TOP_MARGIN + idx * CONFIG.BAR_HEIGHT + 15

        # Troncature des textes
        code = str(row['internal code'])[:CONFIG.MAX_CODE_LENGTH]
        target = str(row['target'])[:CONFIG.MAX_TARGET_LENGTH]

        # Labels de gauche
        parts.append(f'<text x="10" y="{y_pos + 5}" class="gantt-label gantt-label-code">{code}</text>')
        parts.append(f'<text x="10" y="{y_pos + 18}" class="gantt-label gantt-label-target">{target}</text>')

        # Calcul des positions
        start_x = self._compute_x_position(row['start_date'], chart_width)
        end_x = self._compute_x_position(row['end_date'], chart_width)
        bar_width = max(end_x - start_x, CONFIG.MIN_BAR_WIDTH)
        color = self.colors[idx % len(self.colors)]

        # Données pour tooltip
        ho_date = row['start_date'].strftime('%d/%m/%Y')
        val_date = row['end_date'].strftime('%d/%m/%Y')
        duration = int(row['duration_days'])

        # Rectangle de la barre
        tooltip = f"{code} | {target} | HO: {ho_date} -&gt; Val: {val_date} | {duration}j"
        parts.append(
            f'<rect x="{start_x}" y="{y_pos - 8}" width="{bar_width}" height="25" '
            f'fill="{color}" class="gantt-bar" title="{tooltip}"/>'
        )

        # Texte dans la barre (si assez large)
        if bar_width > CONFIG.MIN_BAR_WIDTH_FOR_TEXT:
            text_x = start_x + bar_width / 2
            parts.append(f'<text x="{text_x}" y="{y_pos + 2}" class="gantt-bar-text">{duration}j</text>')

        return parts

    def create_gantt_svg(self, product_line_data: DataFrame, product_line_name: str) -> str:
        """Crée un Gantt chart SVG pour une product line."""
        data = product_line_data.sort_values('start_date').reset_index(drop=True)
        model_count = len(data)

        # Dimensions calculées
        chart_width = CONFIG.TOTAL_WIDTH - CONFIG.LEFT_MARGIN - CONFIG.RIGHT_MARGIN
        total_height = model_count * CONFIG.BAR_HEIGHT + CONFIG.TOP_MARGIN + CONFIG.BOTTOM_MARGIN

        # Construction avec liste
        svg_parts: list[str] = []

        # En-tête
        svg_parts.extend(self._build_svg_header(CONFIG.TOTAL_WIDTH, total_height, product_line_name, model_count))

        # Grille de dates
        svg_parts.extend(self._build_date_grid(chart_width, total_height))

        # Barres Gantt
        for idx, row in data.iterrows():
            svg_parts.extend(self._build_gantt_bar(row, idx, chart_width))

        # Axe horizontal et légende
        axis_y = CONFIG.TOP_MARGIN + model_count * CONFIG.BAR_HEIGHT
        svg_parts.append(
            f'<line x1="{CONFIG.LEFT_MARGIN}" y1="{axis_y}" '
            f'x2="{CONFIG.TOTAL_WIDTH - CONFIG.RIGHT_MARGIN}" y2="{axis_y}" '
            f'stroke="#34495e" stroke-width="2"/>'
        )
        svg_parts.append(
            f'<text x="{CONFIG.LEFT_MARGIN}" y="{total_height - 30}" class="gantt-legend">'
            f'Disponibilite HO → Fin Validation</text>'
        )
        svg_parts.append('</svg>')

        return ''.join(svg_parts)

    def generate_gantt_charts(self) -> 'GanttGenerator':
        """Génère tous les Gantt charts."""
        grouped = self.df_clean.groupby('product line', sort=True)

        for pl_name, pl_data in grouped:
            self.gantt_svgs[pl_name] = self.create_gantt_svg(pl_data, pl_name)

        return self

    def _build_cover_slide(self) -> str:
        """Construit la slide de couverture."""
        return f"""
<div class="slide-container">
    <div class="slide-header">
        <div class="slide-title">Diagrammes de Gantt</div>
        <div class="slide-subtitle">Modeles Murins Genetiquement Modifies - CRUPPE</div>
    </div>
    <div style="text-align: center; margin-top: 60px;">
        <p style="font-size: 18px; color: #7f8c8d; margin: 20px 0;">
            <strong>{len(self.df_clean)} modeles</strong> repartis dans
            <strong>{len(self.gantt_svgs)} product lines</strong>
        </p>
        <p style="font-size: 14px; color: #95a5a6; margin: 20px 0;">
            Periode: {self.min_date.strftime('%d/%m/%Y')} → {self.max_date.strftime('%d/%m/%Y')}
        </p>
        <p style="font-size: 12px; color: #bdc3c7; margin: 40px 0;">
            Genere le {datetime.now().strftime('%d/%m/%Y a %H:%M')}
        </p>
    </div>
</div>"""

    def _build_gantt_slide(self, idx: int, product_line: str, svg: str) -> str:
        """Construit une slide Gantt."""
        return f"""
<div class="slide-container">
    <div class="slide-header">
        <div class="slide-title">Produit {idx}/{len(self.gantt_svgs)}</div>
        <div class="slide-subtitle">{product_line}</div>
    </div>
    <div class="gantt-wrapper">{svg}</div>
</div>"""

    def _build_summary_slide(self) -> str:
        """Construit la slide de résumé."""
        counts = self.df_clean['product line'].value_counts().sort_values(ascending=False)

        rows = ''.join(
            f'<tr><td>{pl}</td><td>{count}</td></tr>'
            for pl, count in counts.items()
        )

        return f"""
<div class="slide-container">
    <div class="slide-header">
        <div class="slide-title">Resume</div>
        <div class="slide-subtitle">Vue d'ensemble des donnees</div>
    </div>
    <div style="margin-top: 40px;">
        <table class="summary-table">
            <tr><th>Product Line</th><th>Modeles</th></tr>
            {rows}
        </table>
    </div>
</div>"""

    def generate_html(self) -> str:
        """Génère le contenu HTML complet."""
        html_parts: list[str] = [
            '<!DOCTYPE html>',
            '<html lang="fr">',
            '<head>',
            '<meta charset="UTF-8">',
            '<meta name="viewport" content="width=device-width, initial-scale=1.0">',
            '<title>Gantt Charts - Modeles Murins CRUPPE</title>',
            f'<style>{HTML_STYLES}</style>',
            '</head>',
            '<body>',
        ]

        html_parts.append(self._build_cover_slide())

        for idx, (product_line, svg) in enumerate(self.gantt_svgs.items(), 1):
            html_parts.append(self._build_gantt_slide(idx, product_line, svg))

        html_parts.append(self._build_summary_slide())
        html_parts.extend(['</body>', '</html>'])

        return ''.join(html_parts)

    def generate_csv(self) -> str:
        """Génère le contenu CSV."""
        column_mapping = {
            'product line': 'Product Line',
            'internal code': 'Internal Code',
            'target': 'Target',
            'status': 'Status',
            'start_date': 'Date Disponibilite HO',
            'end_date': 'Date Fin Validation',
            'duration_days': 'Duree (jours)',
        }

        # Filtrer uniquement les colonnes qui existent
        available_cols = [col for col in column_mapping.keys() if col in self.df_clean.columns]

        export_df = (
            self.df_clean[available_cols]
            .rename(columns={k: v for k, v in column_mapping.items() if k in available_cols})
            .sort_values('Product Line' if 'product line' in available_cols else available_cols[0])
        )

        return export_df.to_csv(index=False, encoding='utf-8')

    def generate_pptx(self) -> bytes:
        """Génère une présentation PowerPoint avec les diagrammes de Gantt."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Slide de couverture
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Titre
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Diagrammes de Gantt"
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Sous-titre
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"Modeles Murins - {len(self.df_clean)} modeles dans {len(self.gantt_svgs)} product lines"
        subtitle_para.font.size = Pt(24)
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(0.5))
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = f"Genere le {datetime.now().strftime('%d/%m/%Y a %H:%M')}"
        date_para.font.size = Pt(14)
        date_para.alignment = PP_ALIGN.CENTER

        # Slides pour chaque product line
        grouped = self.df_clean.groupby('product line', sort=True)

        for pl_name, pl_data in grouped:
            slide = prs.slides.add_slide(slide_layout)

            # Titre de la slide
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = str(pl_name)
            title_para.font.size = Pt(24)
            title_para.font.bold = True

            # Créer le diagramme Gantt avec matplotlib
            fig = create_matplotlib_gantt(pl_data, str(pl_name), self.colors)
            img_bytes = fig_to_image_bytes(fig, 'png', 150)

            # Ajouter l'image
            image_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(image_stream, Inches(0.3), Inches(1), width=Inches(12.7))

        # Slide de résumé
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(0.6))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Resume"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Tableau de résumé
        counts = self.df_clean['product line'].value_counts().sort_values(ascending=False)
        rows = len(counts) + 1
        cols = 2

        table = slide.shapes.add_table(rows, cols, Inches(2), Inches(1.5), Inches(9), Inches(0.4 * rows)).table

        # En-têtes
        table.cell(0, 0).text = "Product Line"
        table.cell(0, 1).text = "Modeles"

        for idx, (pl, count) in enumerate(counts.items(), 1):
            if idx < rows:
                table.cell(idx, 0).text = str(pl)
                table.cell(idx, 1).text = str(count)

        # Sauvegarder en bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer.getvalue()

    def generate_docx(self) -> bytes:
        """Génère un document Word avec les diagrammes de Gantt."""
        doc = Document()

        # Titre principal
        title = doc.add_heading('Diagrammes de Gantt', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Sous-titre
        subtitle = doc.add_paragraph()
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = subtitle.add_run('Modeles Murins Genetiquement Modifies - CRUPPE')
        run.font.size = DocxPt(14)

        # Informations générales
        info = doc.add_paragraph()
        info.alignment = WD_ALIGN_PARAGRAPH.CENTER
        info.add_run(f'{len(self.df_clean)} modeles repartis dans {len(self.gantt_svgs)} product lines\n')
        info.add_run(f'Periode: {self.min_date.strftime("%d/%m/%Y")} - {self.max_date.strftime("%d/%m/%Y")}\n')
        info.add_run(f'Genere le {datetime.now().strftime("%d/%m/%Y a %H:%M")}')

        doc.add_page_break()

        # Diagrammes pour chaque product line
        grouped = self.df_clean.groupby('product line', sort=True)

        for pl_name, pl_data in grouped:
            doc.add_heading(str(pl_name), 1)

            # Créer le diagramme Gantt avec matplotlib
            fig = create_matplotlib_gantt(pl_data, str(pl_name), self.colors)
            img_bytes = fig_to_image_bytes(fig, 'png', 150)

            # Ajouter l'image
            image_stream = io.BytesIO(img_bytes)
            doc.add_picture(image_stream, width=DocxInches(6.5))

            # Tableau des détails
            doc.add_paragraph()
            table = doc.add_table(rows=1, cols=4)
            table.style = 'Table Grid'

            # En-têtes
            header_cells = table.rows[0].cells
            header_cells[0].text = 'Code'
            header_cells[1].text = 'Target'
            header_cells[2].text = 'Date Debut'
            header_cells[3].text = 'Date Fin'

            # Données
            for _, row in pl_data.iterrows():
                row_cells = table.add_row().cells
                row_cells[0].text = str(row['internal code'])[:30]
                row_cells[1].text = str(row['target'])[:30]
                row_cells[2].text = row['start_date'].strftime('%d/%m/%Y')
                row_cells[3].text = row['end_date'].strftime('%d/%m/%Y')

            doc.add_page_break()

        # Page de résumé
        doc.add_heading('Resume', 1)

        counts = self.df_clean['product line'].value_counts().sort_values(ascending=False)

        summary_table = doc.add_table(rows=1, cols=2)
        summary_table.style = 'Table Grid'

        header_cells = summary_table.rows[0].cells
        header_cells[0].text = 'Product Line'
        header_cells[1].text = 'Nombre de Modeles'

        for pl, count in counts.items():
            row_cells = summary_table.add_row().cells
            row_cells[0].text = str(pl)
            row_cells[1].text = str(count)

        # Sauvegarder en bytes
        docx_buffer = io.BytesIO()
        doc.save(docx_buffer)
        docx_buffer.seek(0)

        return docx_buffer.getvalue()


# =============================================================================
# APPLICATION STREAMLIT
# =============================================================================

def main():
    """Application Streamlit principale."""
    st.set_page_config(
        page_title="Gantt Chart Generator",
        page_icon="📊",
        layout="wide",
        initial_sidebar_state="expanded"
    )

    st.title("📊 Generateur de Diagrammes de Gantt")
    st.markdown("*Modeles Murins Genetiquement Modifies - CRUPPE*")

    # Sidebar pour les options
    with st.sidebar:
        st.header("⚙️ Configuration")

        # Upload du fichier
        uploaded_file = st.file_uploader(
            "📁 Charger un fichier Excel",
            type=['xlsx', 'xls'],
            help="Le fichier doit contenir une feuille 'tableau complet'"
        )

        st.divider()

        # Color picker
        st.subheader("🎨 Couleurs des barres")
        st.caption("Personnalisez les couleurs des diagrammes")

        num_colors = st.slider("Nombre de couleurs", 3, 10, 5)

        colors = []
        cols = st.columns(2)
        for i in range(num_colors):
            default_color = CONFIG.DEFAULT_COLORS[i % len(CONFIG.DEFAULT_COLORS)]
            with cols[i % 2]:
                color = st.color_picker(f"Couleur {i+1}", default_color, key=f"color_{i}")
                colors.append(color)

        st.divider()

        # Options d'export
        st.subheader("📤 Formats d'export")
        export_html = st.checkbox("HTML", value=True)
        export_csv = st.checkbox("CSV", value=True)
        export_pptx = st.checkbox("PowerPoint (PPTX)", value=True)
        export_docx = st.checkbox("Word (DOCX)", value=True)

    # Zone principale
    if uploaded_file is not None:
        try:
            # Charger les données
            with st.spinner("Chargement des donnees..."):
                df = pd.read_excel(uploaded_file, sheet_name='tableau complet')
                generator = GanttGenerator(df, colors)
                generator.load_data()
                generator.generate_gantt_charts()

            # Afficher les statistiques
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Modeles valides", len(generator.df_clean))
            with col2:
                st.metric("Product Lines", len(generator.gantt_svgs))
            with col3:
                st.metric("Periode", f"{generator.min_date.strftime('%m/%Y')} - {generator.max_date.strftime('%m/%Y')}")

            st.divider()

            # Tabs pour navigation
            tab1, tab2, tab3 = st.tabs(["📈 Diagrammes", "📋 Donnees", "📥 Telecharger"])

            with tab1:
                st.subheader("Diagrammes de Gantt par Product Line")

                # Sélecteur de product line
                product_lines = list(generator.gantt_svgs.keys())
                selected_pl = st.selectbox("Selectionner une Product Line", product_lines)

                if selected_pl:
                    # Afficher le SVG
                    svg_content = generator.gantt_svgs[selected_pl]
                    st.components.v1.html(svg_content, height=600, scrolling=True)

                    # Afficher aussi avec matplotlib
                    with st.expander("🖼️ Version image (pour export)"):
                        pl_data = generator.df_clean[generator.df_clean['product line'] == selected_pl]
                        fig = create_matplotlib_gantt(pl_data, selected_pl, colors)
                        st.pyplot(fig)

            with tab2:
                st.subheader("Donnees des modeles")

                # Filtres
                col1, col2 = st.columns(2)
                with col1:
                    filter_pl = st.multiselect(
                        "Filtrer par Product Line",
                        options=product_lines,
                        default=[]
                    )

                # Afficher les données
                display_df = generator.df_clean.copy()
                if filter_pl:
                    display_df = display_df[display_df['product line'].isin(filter_pl)]

                # Colonnes à afficher
                display_cols = ['product line', 'internal code', 'target', 'start_date', 'end_date', 'duration_days']
                available_display_cols = [c for c in display_cols if c in display_df.columns]

                st.dataframe(
                    display_df[available_display_cols],
                    use_container_width=True,
                    hide_index=True
                )

                st.caption(f"Total: {len(display_df)} modeles")

            with tab3:
                st.subheader("Telecharger les fichiers")

                col1, col2 = st.columns(2)

                with col1:
                    if export_html:
                        html_content = generator.generate_html()
                        st.download_button(
                            label="📄 Telecharger HTML",
                            data=html_content,
                            file_name="gantt_charts.html",
                            mime="text/html",
                            use_container_width=True
                        )

                    if export_csv:
                        csv_content = generator.generate_csv()
                        st.download_button(
                            label="📊 Telecharger CSV",
                            data=csv_content,
                            file_name="gantt_export.csv",
                            mime="text/csv",
                            use_container_width=True
                        )

                with col2:
                    if export_pptx:
                        with st.spinner("Generation du PowerPoint..."):
                            pptx_content = generator.generate_pptx()
                        st.download_button(
                            label="📽️ Telecharger PPTX",
                            data=pptx_content,
                            file_name="gantt_presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )

                    if export_docx:
                        with st.spinner("Generation du document Word..."):
                            docx_content = generator.generate_docx()
                        st.download_button(
                            label="📝 Telecharger DOCX",
                            data=docx_content,
                            file_name="gantt_document.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            use_container_width=True
                        )

                st.divider()
                st.info("💡 Les fichiers PPTX et DOCX incluent les diagrammes sous forme d'images avec les couleurs personnalisees.")

        except Exception as e:
            st.error(f"❌ Erreur lors du traitement: {str(e)}")
            st.exception(e)

    else:
        # Message d'accueil
        st.info("👆 Chargez un fichier Excel pour commencer")

        with st.expander("ℹ️ Format attendu du fichier"):
            st.markdown("""
            Le fichier Excel doit contenir une feuille nommée **'tableau complet'** avec les colonnes suivantes:

            | Colonne | Description |
            |---------|-------------|
            | `internal code` | Code interne du modele |
            | `target` | Cible du modele |
            | `product line` | Ligne de produit |
            | `date disponibilité HO` | Date de debut |
            | `data de fin de validation...` | Date de fin |
            | `status` | Statut (optionnel) |
            """)

        with st.expander("🎨 Personnalisation des couleurs"):
            st.markdown("""
            Utilisez le panneau de gauche pour:
            - Choisir le nombre de couleurs (3-10)
            - Personnaliser chaque couleur avec le color picker
            - Les couleurs seront appliquees aux barres du diagramme de Gantt
            """)


if __name__ == '__main__':
    main()
